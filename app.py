import base64
import mimetypes
import json
from typing import Dict, List
from fastapi import FastAPI, File, UploadFile, HTTPException
from fastapi.responses import HTMLResponse, JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn
from anthropic import Anthropic
import os
from dotenv import load_dotenv
import io
import requests
from prompts import IMAGE_PROMPT,VIDEO_PROMPT, ANALYZE_TEXT_SEGMENT_PROMPT,ANALYZE_TEXT_CHUNK_PROMPT
from config import SUPPORTED_MIME_TYPES,SUPPORTED_TEXT_TYPES,SUPPORTED_VIDEO_TYPES
# Load environment variables
load_dotenv()

# Initialize FastAPI app
app = FastAPI(title="AI Image Detector")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Mount static files
app.mount("/static", StaticFiles(directory="static"), name="static")

# Anthropic client
ANTHROPIC_API_KEY = os.getenv(
    "ANTHROPIC_API_KEY"
)
client = Anthropic(api_key=ANTHROPIC_API_KEY)


class TextAnalysisRequest(BaseModel):
    text: str

class URLRequest(BaseModel):
    url: str

def chunk_text(text: str, chunk_size: int = 8000, overlap: int = 200) -> List[str]:
    """
    Split text into chunks with overlap to maintain context.
    
    Args:
        text: The text to chunk
        chunk_size: Maximum size of each chunk (in characters)
        overlap: Number of characters to overlap between chunks
    
    Returns:
        List of text chunks
    """
    if len(text) <= chunk_size:
        return [text]
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + chunk_size
        
        # If not the last chunk, try to break at a sentence boundary
        if end < len(text):
            # Look for sentence endings within the last 200 chars
            sentence_endings = ['. ', '.\n', '! ', '!\n', '? ', '?\n']
            for i in range(end, max(start + chunk_size - 200, start), -1):
                if any(text[i:i+2] == ending for ending in sentence_endings):
                    end = i + 2
                    break
        
        chunk = text[start:end]
        chunks.append(chunk.strip())
        
        # Move start position with overlap
        start = end - overlap if end < len(text) else end
    
    return chunks

def split_into_segments(text: str) -> List[Dict[str, any]]:
    """
    Split text into segments (sentences/paragraphs) for detailed analysis.
    
    Args:
        text: The text to segment
    
    Returns:
        List of segments with their text and position
    """
    import re
    
    # For very short text, return as single segment
    if len(text) < 100:
        return [{
            "text": text,
            "start": 0,
            "end": len(text)
        }]
    
    # Split by paragraphs first (double newlines)
    paragraphs = re.split(r'\n\s*\n', text)
    
    segments = []
    current_pos = 0
    
    for para in paragraphs:
        if not para.strip():
            current_pos += len(para) + 2  # +2 for the double newline
            continue
        
        # Split paragraph into sentences
        sentences = re.split(r'([.!?]\s+)', para)
        
        current_segment = ""
        segment_start = current_pos
        
        for i in range(0, len(sentences), 2):
            if i < len(sentences):
                sentence = sentences[i]
                if i + 1 < len(sentences):
                    sentence += sentences[i + 1]
                
                current_segment += sentence
                
                # Create segment every 2-3 sentences or if segment is long enough (200-400 chars)
                if len(current_segment) > 250 or (i + 2 >= len(sentences)):
                    if current_segment.strip():
                        segment_text = current_segment.strip()
                        segments.append({
                            "text": segment_text,
                            "start": segment_start,
                            "end": segment_start + len(segment_text)
                        })
                        segment_start += len(current_segment)
                        current_segment = ""
        
        # Handle remaining text in paragraph
        if current_segment.strip():
            segment_text = current_segment.strip()
            segments.append({
                "text": segment_text,
                "start": segment_start,
                "end": segment_start + len(segment_text)
            })
        
        current_pos += len(para) + 2
    
    # If no segments created, split by sentences only
    if not segments:
        sentences = re.findall(r'[^.!?]+[.!?]\s*', text)
        current_pos = 0
        current_segment = ""
        segment_start = 0
        
        for sentence in sentences:
            current_segment += sentence
            if len(current_segment) > 200:
                if current_segment.strip():
                    segment_text = current_segment.strip()
                    segments.append({
                        "text": segment_text,
                        "start": segment_start,
                        "end": segment_start + len(segment_text)
                    })
                    segment_start += len(current_segment)
                    current_segment = ""
        
        if current_segment.strip():
            segment_text = current_segment.strip()
            segments.append({
                "text": segment_text,
                "start": segment_start,
                "end": segment_start + len(segment_text)
            })
    
    # Final fallback: return whole text as single segment
    if not segments:
        segments.append({
            "text": text,
            "start": 0,
            "end": len(text)
        })
    
    return segments

def get_mime_type(filename: str) -> str:
    """Get MIME type from filename"""
    mime_type, _ = mimetypes.guess_type(filename)
    if mime_type not in SUPPORTED_MIME_TYPES:
        raise ValueError(f"Unsupported image type: {mime_type}")
    return mime_type

@app.get("/", response_class=HTMLResponse)
async def read_root():
    """Serve the main HTML page"""
    # Fixed: Added encoding='utf-8' to handle Unicode characters
    with open("static/index.html", "r", encoding="utf-8") as f:
        return f.read()

@app.post("/analyze")
async def analyze_image(file: UploadFile = File(...)) -> JSONResponse:
    """Analyze uploaded image to determine if it's AI-generated or real"""
    try:
        # Validate file type
        mime_type = get_mime_type(file.filename)

        # Read and encode image
        contents = await file.read()
        image_base64 = base64.b64encode(contents).decode("utf-8")

        # Call Claude API
        response = client.messages.create(
            model="claude-sonnet-4-5-20250929",
            max_tokens=1000,
            temperature=0,
            system="You are an AI and Real image identifier. Respond ONLY with valid JSON.",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_base64
                            }
                        },
                        {
                            "type": "text",
                            "text": IMAGE_PROMPT
                        }
                    ]
                }
            ]
        )

        # Parse response
        response_text = response.content[0].text

        # Try to extract JSON if it's wrapped in markdown code blocks
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()

        # Parse JSON response
        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            # Fallback parsing
            result = {
                "classification": "Analysis Complete",
                "confidence": 75,
                "reasoning": response_text,
                "details": ["Analysis provided above"]
            }

        return JSONResponse(content={
            "success": True,
            "result": result
        })

    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        print(e)
        raise HTTPException(status_code=500, detail=f"Error analyzing image: {str(e)}")

@app.post("/analyze-image-url")
async def analyze_image_url(request: URLRequest) -> JSONResponse:
    """Analyze image from URL or base64 data to determine if it's AI-generated or real"""
    try:
        url = request.url

        # Check if input is base64 data
        if url.startswith('data:image/'):
            # Parse data URI format: data:image/png;base64,iVBORw0KG...
            try:
                # Extract MIME type and base64 data
                header, base64_data = url.split(',', 1)
                mime_type = header.split(';')[0].split(':')[1]

                # Validate MIME type
                if mime_type not in SUPPORTED_MIME_TYPES:
                    raise ValueError(f"Unsupported image type: {mime_type}")

                # Use the base64 data directly
                image_base64 = base64_data

            except (ValueError, IndexError) as e:
                raise ValueError(f"Invalid base64 data URI format: {str(e)}")

        elif url.startswith('base64,') or url.startswith('data:,'):
            # Handle alternative base64 formats
            try:
                if url.startswith('base64,'):
                    image_base64 = url.split('base64,', 1)[1]
                else:
                    image_base64 = url.split('data:,', 1)[1]

                # Try to decode to validate
                base64.b64decode(image_base64)

                # Default to JPEG if no type specified
                mime_type = 'image/jpeg'

            except Exception as e:
                raise ValueError(f"Invalid base64 data: {str(e)}")

        else:
            # Regular URL - download the image
            response = requests.get(url, timeout=30)
            response.raise_for_status()

            contents = response.content

            # Determine MIME type
            content_type = response.headers.get('Content-Type', '')
            if 'image' in content_type:
                mime_type = content_type.split(';')[0]
            else:
                # Try to guess from URL extension
                file_ext = url.split('.')[-1].lower().split('?')[0]
                ext_to_mime = {
                    'jpg': 'image/jpeg',
                    'jpeg': 'image/jpeg',
                    'png': 'image/png',
                    'gif': 'image/gif',
                    'webp': 'image/webp'
                }
                mime_type = ext_to_mime.get(file_ext, 'image/jpeg')

            # Validate MIME type
            if mime_type not in SUPPORTED_MIME_TYPES:
                raise ValueError(f"Unsupported image type: {mime_type}")

            # Encode image
            image_base64 = base64.b64encode(contents).decode("utf-8")

        # Call Claude API
        api_response = client.messages.create(
            model="claude-sonnet-4-5-20250929",
            max_tokens=1000,
            temperature=0,
            system="You are an AI and Real image identifier. Respond ONLY with valid JSON.",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_base64
                            }
                        },
                        {
                            "type": "text",
                            "text": IMAGE_PROMPT
                        }
                    ]
                }
            ]
        )

        # Parse response
        response_text = api_response.content[0].text

        # Try to extract JSON if it's wrapped in markdown code blocks
        if "```json" in response_text:
            response_text = response_text.split("```json")[1].split("```")[0].strip()
        elif "```" in response_text:
            response_text = response_text.split("```")[1].split("```")[0].strip()

        # Parse JSON response
        try:
            result = json.loads(response_text)
        except json.JSONDecodeError:
            # Fallback parsing
            result = {
                "classification": "Analysis Complete",
                "confidence": 75,
                "reasoning": response_text,
                "details": ["Analysis provided above"]
            }

        return JSONResponse(content={
            "success": True,
            "result": result
        })

    except requests.exceptions.RequestException as e:
        raise HTTPException(status_code=400, detail=f"Error downloading image from URL: {str(e)}")
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error analyzing image: {str(e)}")

@app.post("/extract-text")
async def extract_text(file: UploadFile = File(...)) -> JSONResponse:
    """Extract text from uploaded document (PDF, DOCX, PPT, TXT)"""
    try:
        file_ext = file.filename.split('.')[-1].lower() if '.' in file.filename else ''
        contents = await file.read()

        text = ""

        if file_ext == 'txt':
            text = contents.decode('utf-8', errors='ignore')
        elif file_ext == 'pdf':
            try:
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
                text = "\n".join([page.extract_text() for page in pdf_reader.pages])
            except ImportError:
                # Fallback: try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(io.BytesIO(contents)) as pdf:
                        text = "\n".join([page.extract_text() or "" for page in pdf.pages])
                except ImportError:
                    raise HTTPException(status_code=500, detail="PDF extraction library not installed. Please install PyPDF2 or pdfplumber.")
        elif file_ext == 'docx':
            try:
                from docx import Document
                doc = Document(io.BytesIO(contents))
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            except ImportError:
                raise HTTPException(status_code=500, detail="DOCX extraction library not installed. Please install python-docx.")
        elif file_ext in ['ppt', 'pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(contents))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
            except ImportError:
                raise HTTPException(status_code=500, detail="PPT extraction library not installed. Please install python-pptx.")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")

        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the file")

        return JSONResponse(content={
            "success": True,
            "text": text
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text: {str(e)}")

@app.post("/extract-text-url")
async def extract_text_url(request: URLRequest) -> JSONResponse:
    """Extract text from document URL (PDF, DOCX, PPT, TXT)"""
    try:
        url = request.url

        # Convert Google Docs/Drive URLs to direct download links
        if 'docs.google.com' in url or 'drive.google.com' in url:
            # Extract document ID from various Google URL formats
            doc_id = None

            # Google Docs format: docs.google.com/document/d/{ID}/edit
            if 'docs.google.com/document/d/' in url:
                doc_id = url.split('/document/d/')[1].split('/')[0]
                # Export as DOCX
                url = f'https://docs.google.com/document/d/{doc_id}/export?format=docx'
                file_ext = 'docx'
            # Google Sheets format: docs.google.com/spreadsheets/d/{ID}/edit
            elif 'docs.google.com/spreadsheets/d/' in url:
                doc_id = url.split('/spreadsheets/d/')[1].split('/')[0]
                # Export as PDF (easier to extract text from)
                url = f'https://docs.google.com/spreadsheets/d/{doc_id}/export?format=pdf'
                file_ext = 'pdf'
            # Google Slides format: docs.google.com/presentation/d/{ID}/edit
            elif 'docs.google.com/presentation/d/' in url:
                doc_id = url.split('/presentation/d/')[1].split('/')[0]
                # Export as PPTX
                url = f'https://docs.google.com/presentation/d/{doc_id}/export?format=pptx'
                file_ext = 'pptx'
            # Google Drive format: drive.google.com/file/d/{ID}/view
            elif 'drive.google.com/file/d/' in url:
                doc_id = url.split('/file/d/')[1].split('/')[0]
                # Direct download
                url = f'https://drive.google.com/uc?export=download&id={doc_id}'
                file_ext = ''  # Will be determined later
            # Google Drive open format: drive.google.com/open?id={ID}
            elif 'drive.google.com/open?id=' in url:
                doc_id = url.split('id=')[1].split('&')[0]
                url = f'https://drive.google.com/uc?export=download&id={doc_id}'
                file_ext = ''

        # Download the file from URL
        response = requests.get(url, timeout=30, allow_redirects=True)
        response.raise_for_status()

        contents = response.content

        # Determine file type from URL or Content-Type header
        if 'file_ext' not in locals() or not file_ext:
            file_ext = url.split('.')[-1].lower().split('?')[0] if '.' in url else ''

        if not file_ext:
            content_type = response.headers.get('Content-Type', '')
            if 'pdf' in content_type:
                file_ext = 'pdf'
            elif 'word' in content_type or 'docx' in content_type or 'officedocument.wordprocessing' in content_type:
                file_ext = 'docx'
            elif 'powerpoint' in content_type or 'pptx' in content_type or 'officedocument.presentation' in content_type:
                file_ext = 'pptx'
            elif 'text' in content_type:
                file_ext = 'txt'

        text = ""

        if file_ext == 'txt':
            text = contents.decode('utf-8', errors='ignore')
        elif file_ext == 'pdf':
            try:
                import PyPDF2
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(contents))
                text = "\n".join([page.extract_text() for page in pdf_reader.pages])
            except ImportError:
                try:
                    import pdfplumber
                    with pdfplumber.open(io.BytesIO(contents)) as pdf:
                        text = "\n".join([page.extract_text() or "" for page in pdf.pages])
                except ImportError:
                    raise HTTPException(status_code=500, detail="PDF extraction library not installed. Please install PyPDF2 or pdfplumber.")
        elif file_ext == 'docx':
            try:
                from docx import Document
                doc = Document(io.BytesIO(contents))
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            except ImportError:
                raise HTTPException(status_code=500, detail="DOCX extraction library not installed. Please install python-docx.")
        elif file_ext in ['ppt', 'pptx']:
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(contents))
                text_parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)
                text = "\n".join(text_parts)
            except ImportError:
                raise HTTPException(status_code=500, detail="PPT extraction library not installed. Please install python-pptx.")
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_ext}")

        if not text.strip():
            raise HTTPException(status_code=400, detail="No text could be extracted from the file")

        return JSONResponse(content={
            "success": True,
            "text": text
        })

    except requests.exceptions.RequestException as e:
        raise HTTPException(status_code=400, detail=f"Error downloading file from URL: {str(e)}")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error extracting text: {str(e)}")

async def analyze_text_segment(segment: str) -> Dict:
    """Analyze a single text segment"""
    response = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=800,
        temperature=0,
        system="You are an AI text detector. Analyze text segments to determine if they're AI-generated or human-written. Respond ONLY with valid JSON.",
        messages=[
            {
                "role": "user",
                "content": ANALYZE_TEXT_SEGMENT_PROMPT.format(segment=segment)
            }
        ]
    )
    
    response_text = response.content[0].text
    
    # Try to extract JSON if it's wrapped in markdown code blocks
    if "```json" in response_text:
        response_text = response_text.split("```json")[1].split("```")[0].strip()
    elif "```" in response_text:
        response_text = response_text.split("```")[1].split("```")[0].strip()
    
    try:
        return json.loads(response_text)
    except json.JSONDecodeError:
        return {
            "classification": "Human-Written",
            "confidence": 50
        }

async def analyze_text_chunk(chunk: str, chunk_num: int, total_chunks: int) -> Dict:
    """Analyze a single chunk of text"""
    response = client.messages.create(
        model="claude-sonnet-4-5-20250929",
        max_tokens=1000,
        temperature=0,
        system="You are an AI text detector. Analyze text to determine if it's AI-generated or human-written. Respond ONLY with valid JSON.",
        messages=[
            {
                "role": "user",
                "content": ANALYZE_TEXT_CHUNK_PROMPT.format(chunk_num=chunk_num,total_chunks=total_chunks,chunk=chunk)
            }
        ]
    )
    
    response_text = response.content[0].text
    
    # Try to extract JSON if it's wrapped in markdown code blocks
    if "```json" in response_text:
        response_text = response_text.split("```json")[1].split("```")[0].strip()
    elif "```" in response_text:
        response_text = response_text.split("```")[1].split("```")[0].strip()
    
    try:
        return json.loads(response_text)
    except json.JSONDecodeError:
        return {
            "classification": "Analysis Complete",
            "confidence": 75,
            "reasoning": response_text,
            "details": ["Analysis provided above"]
        }

def aggregate_chunk_results(chunk_results: List[Dict]) -> Dict:
    """Aggregate results from multiple chunks into a single result"""
    if not chunk_results:
        return {
            "classification": "Analysis Complete",
            "confidence": 50,
            "reasoning": "No chunks analyzed",
            "details": []
        }
    
    # Count classifications
    ai_count = sum(1 for r in chunk_results if "ai" in r.get("classification", "").lower())
    human_count = len(chunk_results) - ai_count
    
    # Calculate average confidence
    avg_confidence = sum(r.get("confidence", 50) for r in chunk_results) / len(chunk_results)
    
    # Determine overall classification
    if ai_count > human_count:
        classification = "AI-Generated"
        confidence = int((ai_count / len(chunk_results)) * avg_confidence)
    elif human_count > ai_count:
        classification = "Human-Written"
        confidence = int((human_count / len(chunk_results)) * avg_confidence)
    else:
        classification = "Mixed/Uncertain"
        confidence = int(avg_confidence)
    
    # Aggregate reasoning and details
    all_reasonings = [r.get("reasoning", "") for r in chunk_results if r.get("reasoning")]
    all_details = []
    for r in chunk_results:
        all_details.extend(r.get("details", []))
    
    # Remove duplicates and limit details
    unique_details = list(dict.fromkeys(all_details))[:5]  # Keep first 5 unique details
    
    reasoning = f"Analyzed {len(chunk_results)} chunk(s). " + " ".join(all_reasonings[:2])  # Combine first 2 reasonings
    
    return {
        "classification": classification,
        "confidence": min(100, max(0, confidence)),  # Ensure 0-100 range
        "reasoning": reasoning,
        "details": unique_details
    }

@app.post("/analyze-text")
async def analyze_text(request: TextAnalysisRequest) -> JSONResponse:
    """Analyze text to determine if it's AI-generated or human-written with segment-level highlighting"""
    try:
        if not request.text or not request.text.strip():
            raise HTTPException(status_code=400, detail="Text cannot be empty")
        
        text = request.text.strip()
        
        # Split text into segments for detailed analysis
        segments = split_into_segments(text)
        
        # Analyze each segment
        segment_results = []
        for segment_data in segments:
            segment_text = segment_data["text"]
            segment_analysis = await analyze_text_segment(segment_text)
            segment_results.append({
                "text": segment_text,
                "start": segment_data["start"],
                "end": segment_data["end"],
                "classification": segment_analysis.get("classification", "Human-Written"),
                "confidence": segment_analysis.get("confidence", 50)
            })
        
        # Calculate percentages based on text length
        total_length = len(text)
        ai_generated_length = 0
        slightly_ai_length = 0
        human_written_length = 0
        
        for segment in segment_results:
            segment_length = segment["end"] - segment["start"]
            classification = segment["classification"].lower()
            
            if "ai-generated" in classification or "ai generated" in classification:
                ai_generated_length += segment_length
            elif "slightly" in classification or "slightly-ai" in classification:
                slightly_ai_length += segment_length
            else:
                human_written_length += segment_length
        
        # Calculate percentages
        ai_percentage = round((ai_generated_length / total_length * 100), 2) if total_length > 0 else 0
        slightly_ai_percentage = round((slightly_ai_length / total_length * 100), 2) if total_length > 0 else 0
        human_percentage = round((human_written_length / total_length * 100), 2) if total_length > 0 else 0
        
        # Total AI percentage (AI + Slightly AI)
        total_ai_percentage = round(ai_percentage + slightly_ai_percentage, 2)
        
        # Also get overall analysis for summary
        chunks = chunk_text(text, chunk_size=8000, overlap=200)
        
        if len(chunks) == 1:
            overall_result = await analyze_text_chunk(chunks[0], 1, 1)
        else:
            chunk_results = []
            for i, chunk in enumerate(chunks, 1):
                chunk_result = await analyze_text_chunk(chunk, i, len(chunks))
                chunk_results.append(chunk_result)
            overall_result = aggregate_chunk_results(chunk_results)
        
        return JSONResponse(content={
            "success": True,
            "result": overall_result,
            "original_text": text,
            "segments": segment_results,
            "statistics": {
                "total_length": total_length,
                "ai_generated_length": ai_generated_length,
                "slightly_ai_length": slightly_ai_length,
                "human_written_length": human_written_length,
                "ai_percentage": ai_percentage,
                "slightly_ai_percentage": slightly_ai_percentage,
                "human_percentage": human_percentage,
                "total_ai_percentage": total_ai_percentage
            }
        })
        
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error analyzing text: {str(e)}")

def extract_video_frames(video_bytes: bytes, max_frames: int = 5) -> List[bytes]:
    """
    Extract frames from video file.

    Args:
        video_bytes: Video file bytes
        max_frames: Maximum number of frames to extract

    Returns:
        List of frame images as bytes
    """
    try:
        import cv2
        import numpy as np
    except ImportError:
        raise HTTPException(status_code=500, detail="OpenCV library not installed. Please install opencv-python.")

    # Save video bytes to temporary file
    import tempfile
    import os

    # Try to preserve original file extension if possible, otherwise use .mp4
    with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp_file:
        tmp_file.write(video_bytes)
        tmp_path = tmp_file.name

    try:
        # Open video
        cap = cv2.VideoCapture(tmp_path)
        if not cap.isOpened():
            raise ValueError("Could not open video file. The video format may be unsupported or corrupted.")

        # Get video properties
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = cap.get(cv2.CAP_PROP_FPS)

        if total_frames <= 0:
            cap.release()
            raise ValueError("Video file has no frames. The file may be corrupted or in an unsupported format.")

        # Calculate frame interval to sample evenly
        if total_frames <= max_frames:
            frame_indices = list(range(total_frames))
        else:
            frame_indices = [int(i * total_frames / max_frames) for i in range(max_frames)]

        frames = []
        for frame_idx in frame_indices:
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_idx)
            ret, frame = cap.read()
            if ret and frame is not None:
                # Convert frame to JPEG bytes
                success, buffer = cv2.imencode('.jpg', frame)
                if success:
                    frames.append(buffer.tobytes())

        cap.release()

        if not frames:
            raise ValueError("Could not extract any frames from the video. Please try a different video file.")

        return frames

    except Exception as e:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
        raise

    finally:
        # Clean up temporary file
        if os.path.exists(tmp_path):
            try:
                os.unlink(tmp_path)
            except:
                pass

@app.post("/analyze-video")
async def analyze_video(file: UploadFile = File(...)) -> JSONResponse:
    """Analyze uploaded video to determine if it's AI-generated or real"""
    try:
        # Validate file type
        file_ext = file.filename.split('.')[-1].lower() if '.' in file.filename else ''
        contents = await file.read()
        
        # Check if it's a video file
        if file_ext not in ['mp4', 'mov', 'avi', 'webm', 'mkv']:
            raise HTTPException(status_code=400, detail=f"Unsupported video format: {file_ext}")
        
        # Extract frames from video
        frames = extract_video_frames(contents, max_frames=5)
        
        if not frames:
            raise HTTPException(status_code=400, detail="Could not extract frames from video")
        
        # Analyze each frame
        frame_results = []
        for i, frame_bytes in enumerate(frames):
            # Encode frame to base64
            frame_base64 = base64.b64encode(frame_bytes).decode("utf-8")
            
            # Analyze frame using Claude API
            response = client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=1000,
                temperature=0,
                system="You are an AI and Real video frame identifier. Respond ONLY with valid JSON.",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/jpeg",
                                    "data": frame_base64
                                }
                            },
                            {
                                "type": "text",
                                "text": VIDEO_PROMPT.format(frame_number=i+1,frame_length=len(frames))
                            }
                        ]
                    }
                ]
            )
            
            # Parse response
            response_text = response.content[0].text
            
            # Try to extract JSON if it's wrapped in markdown code blocks
            if "```json" in response_text:
                response_text = response_text.split("```json")[1].split("```")[0].strip()
            elif "```" in response_text:
                response_text = response_text.split("```")[1].split("```")[0].strip()
            
            try:
                frame_result = json.loads(response_text)
                frame_results.append(frame_result)
            except json.JSONDecodeError:
                frame_results.append({
                    "classification": "Analysis Complete",
                    "confidence": 75,
                    "reasoning": response_text
                })
        
        # Aggregate frame results
        ai_count = sum(1 for r in frame_results if "ai" in r.get("classification", "").lower())
        real_count = len(frame_results) - ai_count

        # Calculate average confidence
        avg_confidence = sum(r.get("confidence", 50) for r in frame_results) / len(frame_results)

        # Determine overall classification and confidence
        if ai_count > 0:
            # If ANY frame is AI-generated, classify as AI-Generated
            classification = f"{int((ai_count / len(frame_results)) * 100)}% AI Generated"
            confidence = int(avg_confidence)
        else:
            # All frames are real
            classification = "0% AI Generated"
            confidence = int(avg_confidence)

        # Create key observations list (4-5 points)
        key_observations = [
            f"Analyzed {len(frame_results)} frames from the video",
            f"{ai_count} AI-generated frames detected" if ai_count > 0 else "No AI-generated frames detected",
            f"{real_count} real video frames detected"
        ]

        # Add specific observations from frame analyses
        for r in frame_results[:2]:  # Get observations from first 2 frames
            if r.get("reasoning"):
                key_observations.append(r.get("reasoning")[:100])  # Limit length

        # Limit to 5 observations
        key_observations = key_observations[:5]

        result = {
            "classification": classification,
            "confidence": min(100, max(0, confidence)),
            "reasoning": f"Video analysis based on {len(frame_results)} sampled frames",
            "details": key_observations
        }

        return JSONResponse(content={
            "success": True,
            "result": result
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error analyzing video: {str(e)}")

@app.post("/analyze-video-url")
async def analyze_video_url(request: URLRequest) -> JSONResponse:
    """Analyze video from YouTube URL to determine if it's AI-generated or real"""
    try:
        url = request.url

        # Download video from YouTube using yt-dlp
        try:
            import yt_dlp
        except ImportError:
            raise HTTPException(status_code=500, detail="yt-dlp library not installed. Please install yt-dlp.")

        import tempfile
        import os

        # Create temporary directory
        with tempfile.TemporaryDirectory() as tmp_dir:
            output_path = os.path.join(tmp_dir, 'video.mp4')

            # Configure yt-dlp options
            ydl_opts = {
                'format': 'worst[ext=mp4]/worst',  # Use worst quality to avoid bot detection and reduce download size
                'outtmpl': output_path,
                'quiet': True,
                'no_warnings': True,
                'nocheckcertificate': True,
                'user_agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                'referer': 'https://www.youtube.com/',
                'socket_timeout': 30,
                'retries': 3,
                'fragment_retries': 3,
                'extractor_retries': 3,
                'file_access_retries': 3,
                # Server-friendly settings
                'age_limit': None,
                'http_headers': {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
                    'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
                    'Accept-Language': 'en-us,en;q=0.5',
                    'Sec-Fetch-Mode': 'navigate',
                },
            }

            # Try to use browser cookies if available (for local development)
            cookies_file = os.getenv('YOUTUBE_COOKIES_FILE')
            if cookies_file and os.path.exists(cookies_file):
                ydl_opts['cookiefile'] = cookies_file
            else:
                # Try browser cookies with fallback
                try:
                    # Check if running in a browser environment (local development)
                    if os.path.exists(os.path.expanduser('~/.config/google-chrome')) or \
                       os.path.exists(os.path.expanduser('~/Library/Application Support/Google/Chrome')) or \
                       os.path.exists(os.path.expanduser(r'~\AppData\Local\Google\Chrome')):
                        ydl_opts['cookiesfrombrowser'] = ('chrome',)
                except:
                    pass  # Silently fail and continue without cookies

            # Download video
            try:
                with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                    ydl.download([url])
            except Exception as e:
                error_msg = str(e)

                # Check if it's a bot detection error
                if "Sign in to confirm" in error_msg or "bot" in error_msg.lower():
                    raise HTTPException(
                        status_code=403,
                        detail=(
                            "YouTube requires authentication. Please set up cookies:\n\n"
                            "1. Install browser extension 'Get cookies.txt LOCALLY'\n"
                            "2. Visit YouTube and login\n"
                            "3. Export cookies to a file (cookies.txt)\n"
                            "4. Set environment variable: YOUTUBE_COOKIES_FILE=/path/to/cookies.txt\n"
                            "5. Restart the server\n\n"
                            "See README.md for detailed instructions."
                        )
                    )
                else:
                    raise HTTPException(status_code=400, detail=f"Error downloading video from YouTube: {error_msg}")

            # Read the downloaded video
            if not os.path.exists(output_path):
                raise HTTPException(status_code=400, detail="Video download failed")

            with open(output_path, 'rb') as f:
                video_bytes = f.read()

        # Extract frames from video
        frames = extract_video_frames(video_bytes, max_frames=5)

        if not frames:
            raise HTTPException(status_code=400, detail="Could not extract frames from video")

        # Analyze each frame
        frame_results = []
        for i, frame_bytes in enumerate(frames):
            # Encode frame to base64
            frame_base64 = base64.b64encode(frame_bytes).decode("utf-8")

            # Analyze frame using Claude API
            response = client.messages.create(
                model="claude-sonnet-4-5-20250929",
                max_tokens=1000,
                temperature=0,
                system="You are an AI and Real video frame identifier. Respond ONLY with valid JSON.",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/jpeg",
                                    "data": frame_base64
                                }
                            },
                            {
                                "type": "text",
                                "text": VIDEO_PROMPT.format(frame_number=i+1, frame_length=len(frames))
                            }
                        ]
                    }
                ]
            )

            # Parse response
            response_text = response.content[0].text

            # Try to extract JSON if it's wrapped in markdown code blocks
            if "```json" in response_text:
                response_text = response_text.split("```json")[1].split("```")[0].strip()
            elif "```" in response_text:
                response_text = response_text.split("```")[1].split("```")[0].strip()

            try:
                frame_result = json.loads(response_text)
                frame_results.append(frame_result)
            except json.JSONDecodeError:
                frame_results.append({
                    "classification": "Analysis Complete",
                    "confidence": 75,
                    "reasoning": response_text
                })

        # Aggregate frame results
        ai_count = sum(1 for r in frame_results if "ai" in r.get("classification", "").lower())
        real_count = len(frame_results) - ai_count

        # Calculate average confidence
        avg_confidence = sum(r.get("confidence", 50) for r in frame_results) / len(frame_results)

        # Determine overall classification and confidence
        if ai_count > 0:
            # If ANY frame is AI-generated, classify as AI-Generated
            classification = f"{int((ai_count / len(frame_results)) * 100)}% AI Generated"
            confidence = int(avg_confidence)
        else:
            # All frames are real
            classification = "0% AI Generated"
            confidence = int(avg_confidence)

        # Create key observations list (4-5 points)
        key_observations = [
            f"Analyzed {len(frame_results)} frames from the YouTube video",
            f"{ai_count} AI-generated frames detected" if ai_count > 0 else "No AI-generated frames detected",
            f"{real_count} real video frames detected"
        ]

        # Add specific observations from frame analyses
        for r in frame_results[:2]:  # Get observations from first 2 frames
            if r.get("reasoning"):
                key_observations.append(r.get("reasoning")[:100])  # Limit length

        # Limit to 5 observations
        key_observations = key_observations[:5]

        result = {
            "classification": classification,
            "confidence": min(100, max(0, confidence)),
            "reasoning": f"Video analysis based on {len(frame_results)} sampled frames",
            "details": key_observations
        }

        return JSONResponse(content={
            "success": True,
            "result": result
        })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error analyzing video: {str(e)}")

@app.get("/health")
async def health_check():
    """Health check endpoint"""
    return {"status": "healthy"}

if __name__ == "__main__":
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True)